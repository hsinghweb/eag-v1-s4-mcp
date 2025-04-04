# basic import 
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from mcp import types
from PIL import Image as PILImage
import math
import sys
from pywinauto.application import Application
import win32gui
import win32con
import time
import asyncio
from win32api import GetSystemMetrics
from pptx import Presentation
from pptx.util import Inches
import os
from pptx.dml.color import RGBColor
from pptx.util import Pt

# Set up logging
import logging

logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

# instantiate an MCP server client with web visualization enabled
mcp = FastMCP("Calculator", host="127.0.0.1", port=8000, enable_web_ui=True)

# Configure logging for tool operations
@mcp.tool()
def log_tool_start(tool_name, **kwargs):
    logging.info(f"Tool Execution Started - Tool: {tool_name}, Args: {kwargs}")

@mcp.tool()
def log_tool_end(tool_name, result):
    logging.info(f"Tool Execution Completed - Tool: {tool_name}, Result: {result}")

@mcp.tool()
def log_tool_error(tool_name, error):
    logging.error(f"Tool Error Occurred - Tool: {tool_name}, Error: {error}")

# DEFINE TOOLS

#addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    logging.debug(f"Adding numbers: {a} + {b}")
    result = int(a + b)
    logging.info(f"Addition result: {result}")
    return result

@mcp.tool()
def add_list(l: list) -> int:
    """Add all numbers in a list"""
    logging.debug(f"Adding numbers in list: {l}")
    result = sum(l)
    logging.info(f"List sum result: {result}")
    return result

# subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    logging.debug(f"Subtracting numbers: {a} - {b}")
    result = int(a - b)
    logging.info(f"Subtraction result: {result}")
    return result

# multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    logging.debug(f"Multiplying numbers: {a} * {b}")
    result = int(a * b)
    logging.info(f"Multiplication result: {result}")
    return result

#  division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    logging.debug(f"Dividing numbers: {a} / {b}")
    try:
        result = float(a / b)
        logging.info(f"Division result: {result}")
        return result
    except ZeroDivisionError as e:
        logging.error(f"Division by zero error: {e}")
        raise

# power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    logging.debug(f"Calculating power: {a} ** {b}")
    result = int(a ** b)
    logging.info(f"Power operation result: {result}")
    return result

# square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    logging.debug(f"Calculating square root of: {a}")
    try:
        result = float(a ** 0.5)
        logging.info(f"Square root result: {result}")
        return result
    except ValueError as e:
        logging.error(f"Square root error: {e}")
        raise

# cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    logging.debug(f"Calculating cube root of: {a}")
    try:
        result = float(a ** (1/3))
        logging.info(f"Cube root result: {result}")
        return result
    except ValueError as e:
        logging.error(f"Cube root error: {e}")
        raise

# factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """factorial of a number"""
    logging.debug(f"Calculating factorial of: {a}")
    try:
        result = int(math.factorial(a))
        logging.info(f"Factorial result: {result}")
        return result
    except ValueError as e:
        logging.error(f"Factorial error: {e}")
        raise

# log tool
@mcp.tool()
def log(a: int) -> float:
    """log of a number"""
    logging.debug(f"Calculating natural log of: {a}")
    try:
        result = float(math.log(a))
        logging.info(f"Natural log result: {result}")
        return result
    except ValueError as e:
        logging.error(f"Natural log error: {e}")
        raise

# remainder tool
@mcp.tool()
def remainder(a: int, b: int) -> int:
    """remainder of two numbers divison"""
    logging.debug(f"Calculating remainder: {a} % {b}")
    try:
        result = int(a % b)
        logging.info(f"Remainder result: {result}")
        return result
    except ZeroDivisionError as e:
        logging.error(f"Remainder error: {e}")
        raise

# sin tool
@mcp.tool()
def sin(a: int) -> float:
    """sin of a number"""
    logging.debug(f"Calculating sine of: {a}")
    result = float(math.sin(a))
    logging.info(f"Sine result: {result}")
    return result

# cos tool
@mcp.tool()
def cos(a: int) -> float:
    """cos of a number"""
    logging.debug(f"Calculating cosine of: {a}")
    result = float(math.cos(a))
    logging.info(f"Cosine result: {result}")
    return result

# tan tool
@mcp.tool()
def tan(a: int) -> float:
    """tan of a number"""
    logging.debug(f"Calculating tangent of: {a}")
    result = float(math.tan(a))
    logging.info(f"Tangent result: {result}")
    return result

# mine tool
@mcp.tool()
def mine(a: int, b: int) -> int:
    """special mining tool"""
    logging.debug(f"Executing mining operation with: {a}, {b}")
    result = int(a - b - b)
    logging.info(f"Mining operation result: {result}")
    return result

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    print("CALLED: create_thumbnail(image_path: str) -> Image:")
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    return Image(data=img.tobytes(), format="png")

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    print("CALLED: strings_to_chars_to_int(string: str) -> list[int]:")
    return [int(ord(char)) for char in string]

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    print("CALLED: int_list_to_exponential_sum(int_list: list) -> float:")
    return sum(math.exp(i) for i in int_list)

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    print("CALLED: fibonacci_numbers(n: int) -> list:")
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    return fib_sequence[:n]

@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        # Try multiple window title variations
        titles = ["presentation.pptx - PowerPoint", "PowerPoint", "Microsoft PowerPoint"]
        closed = False
        
        for title in titles:
            hwnd = win32gui.FindWindow(None, title)
            if hwnd:
                try:
                    win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
                    closed = True
                except Exception as e:
                    logging.error(f"Failed to close window '{title}': {e}")
                    continue
        
        if closed:
            await asyncio.sleep(1)
            return {"content": [TextContent(type="text", text="PowerPoint closed successfully")]}
        else:
            return {"content": [TextContent(type="text", text="No PowerPoint windows found")]}
    except Exception as e:
        logging.error(f"Error in close_powerpoint: {e}")
        return {"content": [TextContent(type="text", text=f"Error closing PowerPoint: {str(e)}")]}

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        # Close any existing PowerPoint instances first
        await close_powerpoint()
        await asyncio.sleep(2)  # Wait for PowerPoint to close
        
        # Create a new presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Save the presentation
        filename = 'presentation.pptx'
        prs.save(filename)
        await asyncio.sleep(2)  # Wait for file to be saved
        
        # Open the presentation
        os.startfile(filename)
        await asyncio.sleep(3)  # Wait for PowerPoint to open
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation"
                )
            ]
        }
    except Exception as e:
        print(f"Error in open_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(x1: int, y1: int, x2: int, y2: int) -> dict:
    """Draw a rectangle in the first slide of PowerPoint"""
    try:
        # Convert and validate parameters
        try:
            x1, y1, x2, y2 = map(lambda x: int(float(str(x))), [x1, y1, x2, y2])
        except (ValueError, TypeError) as e:
            error_msg = f"Failed to convert parameters to integers: {str(e)}"
            logging.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}

        # Validate coordinates
        if not all(1 <= coord <= 8 for coord in [x1, y1, x2, y2]):
            error_msg = f"Coordinates must be between 1 and 8, got: ({x1},{y1}) to ({x2},{y2})"
            logging.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        if x2 <= x1 or y2 <= y1:
            error_msg = f"End coordinates must be greater than start coordinates"
            logging.error(error_msg)
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        # Close PowerPoint before modifying
        await close_powerpoint()
        await asyncio.sleep(1)
        
        try:
            # Open and modify presentation
            prs = Presentation('presentation.pptx')
            slide = prs.slides[0]
            
            # Convert coordinates to inches
            left = Inches(x1)
            top = Inches(y1)
            width = Inches(x2 - x1)
            height = Inches(y2 - y1)
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                left, top, width, height
            )
            
            # Style the rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(4)
            
            # Save and reopen
            prs.save('presentation.pptx')
            await asyncio.sleep(1)
            os.startfile('presentation.pptx')
            await asyncio.sleep(2)
            
            print("DEBUG: Rectangle drawn successfully")
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Rectangle drawn successfully from ({x1},{y1}) to ({x2},{y2})"
                    )
                ]
            }
            
        except Exception as e:
            error_msg = f"PowerPoint operation failed: {str(e)}"
            print(f"DEBUG: {error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        print(f"DEBUG: {error_msg}")
        print(f"DEBUG: Error type: {type(e)}")
        import traceback
        traceback.print_exc()
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(text: str) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        print(f"DEBUG: Received text to add: {text}")
        print(f"DEBUG: Text type: {type(text)}")
        print(f"DEBUG: Text length: {len(text)}")
        print(f"DEBUG: Text contains newlines: {repr(text)}")
        
        # Wait before adding text
        await asyncio.sleep(5)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        await asyncio.sleep(5)
        
        # Open the existing presentation
        prs = Presentation('presentation.pptx')
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        # Match the rectangle position from draw_rectangle
        left = Inches(2.2)  # Slightly more than rectangle left for margin
        top = Inches(2.5)   # Centered vertically in rectangle
        width = Inches(4.6) # Slightly less than rectangle width for margin
        height = Inches(2)  # Enough height for text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True  # Enable word wrap
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Split text into lines
        lines = text.split('\n')
        print(f"DEBUG: Number of lines: {len(lines)}")
        print(f"DEBUG: Lines to add: {lines}")
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            if line.strip():  # Only add non-empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.alignment = 1  # Center align the text
                
                # Format the text
                run = p.runs[0]
                if "Final Result:" in line:
                    run.font.size = Pt(32)  # Header size
                    run.font.bold = True
                else:
                    run.font.size = Pt(28)  # Value size
                    run.font.bold = True
                
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                p.space_after = Pt(12)  # Add spacing between lines
        
        # Save and wait
        prs.save('presentation.pptx')
        await asyncio.sleep(5)
        
        # Reopen PowerPoint
        os.startfile('presentation.pptx')
        await asyncio.sleep(10)
        
        print(f"DEBUG: Text added successfully: {text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {text}"
                )
            ]
        }
    except Exception as e:
        print(f"Error in add_text_in_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

# DEFINE RESOURCES

# Add a dynamic greeting resource
@mcp.resource("greeting://{name}")
def get_greeting(name: str) -> str:
    """Get a personalized greeting"""
    print("CALLED: get_greeting(name: str) -> str:")
    return f"Hello, {name}!"


# DEFINE AVAILABLE PROMPTS
@mcp.prompt()
def review_code(code: str) -> str:
    return f"Please review this code:\n\n{code}"
    print("CALLED: review_code(code: str) -> str:")


@mcp.prompt()
def debug_error(error: str) -> list[base.Message]:
    return [
        base.UserMessage("I'm seeing this error:"),
        base.UserMessage(error),
        base.AssistantMessage("I'll help debug that. What have you tried so far?"),
    ]

if __name__ == "__main__":
    # Check if running with mcp dev command
    print("STARTING THE SERVER")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        mcp.run()  # Run without transport for dev server
    else:
        mcp.run(transport="stdio")  # Run with stdio for direct execution
